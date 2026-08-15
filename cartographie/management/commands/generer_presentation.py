# -*- coding: utf-8 -*-
"""Génère la Présentation PowerPoint MUKMAP (16:9, animations) — édition MUKESHABA.

Usage :
    python manage.py generer_presentation
"""
import os

from django.conf import settings
from django.core.management.base import BaseCommand

from cartographie.guide_contenu import META


BLEU_FONCE = (0x0F, 0x35, 0x54)
BLEU = (0x1F, 0x4E, 0x79)
BLEU_CLAIR = (0x2E, 0x74, 0xB5)
ROUGE = (0xC0, 0x50, 0x4D)
VERT = (0x2E, 0x8B, 0x57)
BLANC = (0xFF, 0xFF, 0xFF)
GRIS = (0xE8, 0xEC, 0xF3)
GRIS_TEXTE = (0x5A, 0x64, 0x77)
FOND = (0xF4, 0xF7, 0xFB)


class Command(BaseCommand):
    help = "Génère la présentation PowerPoint MUKMAP (16:9, animations) — édition MUKESHABA."

    def handle(self, *args, **options):
        dossier = os.path.join(settings.BASE_DIR, 'cartographie', 'static', 'docs')
        os.makedirs(dossier, exist_ok=True)
        chemin = os.path.join(dossier, 'MUKMAP_Presentation_MUKESHABA.pptx')
        self._construire(chemin)
        self.stdout.write(self.style.SUCCESS("Présentation générée : %s" % chemin))
        return chemin

    # =====================================================================
    def _construire(self, chemin):
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.oxml.ns import qn

        from lxml import etree

        dossier = os.path.join(settings.BASE_DIR, 'cartographie', 'static', 'docs')
        captures = os.path.join(dossier, 'captures')
        logo_app = os.path.join(settings.BASE_DIR, 'cartographie', 'static',
                                'logo', 'MUKMAP.png')
        logo_editeur = os.path.join(settings.BASE_DIR, 'cartographie', 'static',
                                    'logo', 'MUKESHABA.png')

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        BLANK = prs.slide_layouts[6]

        def rgb(t):
            return RGBColor(*t)

        def fond(slide, couleur):
            f = slide.background.fill
            f.solid()
            f.fore_color.rgb = rgb(couleur)

        def boite(slide, x, y, w, h, couleur=None, contour=None, arrondi=False):
            forme = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE if arrondi else MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(h))
            if couleur is not None:
                forme.fill.solid()
                forme.fill.fore_color.rgb = rgb(couleur)
            else:
                forme.fill.background()
            if contour is not None:
                forme.line.color.rgb = rgb(contour)
                forme.line.width = Pt(1)
            else:
                forme.line.fill.background()
            forme.shadow.inherit = False
            return forme

        def texte(slide, x, y, w, h, lignes, taille=18, couleur=GRIS_TEXTE,
                  gras=False, align=PP_ALIGN.LEFT, ancrage=MSO_ANCHOR.TOP,
                  italic=False, police='Calibri'):
            tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = ancrage
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)
            if isinstance(lignes, str):
                lignes = [lignes]
            for i, l in enumerate(lignes):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = align
                if isinstance(l, tuple):
                    txt, opts = l
                else:
                    txt, opts = l, {}
                run = p.add_run()
                run.text = txt
                run.font.size = Pt(opts.get('taille', taille))
                run.font.bold = opts.get('gras', gras)
                run.font.italic = opts.get('italic', italic)
                run.font.name = opts.get('police', police)
                run.font.color.rgb = rgb(opts.get('couleur', couleur))
            return tb

        def puces(slide, x, y, w, h, items, taille=16, couleur=GRIS_TEXTE):
            tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)
            for i, it in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(8)
                r1 = p.add_run()
                r1.text = "•  "
                r1.font.size = Pt(taille)
                r1.font.bold = True
                r1.font.color.rgb = rgb(BLEU_CLAIR)
                if isinstance(it, tuple):
                    t, opts = it
                else:
                    t, opts = it, {}
                r2 = p.add_run()
                r2.text = t
                r2.font.size = Pt(opts.get('taille', taille))
                r2.font.bold = opts.get('gras', False)
                r2.font.color.rgb = rgb(opts.get('couleur', couleur))
                r2.font.name = 'Calibri'
            return tb

        def image(slide, chemin_img, x, y, w=None, h=None, bordure=None):
            pic = slide.shapes.add_picture(chemin_img, Inches(x), Inches(y),
                                           Inches(w) if w else None,
                                           Inches(h) if h else None)
            if bordure is not None:
                pic.line.color.rgb = rgb(bordure)
                pic.line.width = Pt(1.5)
            return pic

        def entete(slide, titre, numero, sous_titre=None):
            boite(slide, 0, 0, 13.333, 0.12, couleur=BLEU_FONCE)
            boite(slide, 0, 0.12, 0.12, 7.5, couleur=BLEU_FONCE)
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(10.5), Inches(1.0))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(0)
            tf.margin_top = Emu(0)
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = titre
            r.font.size = Pt(30)
            r.font.bold = True
            r.font.color.rgb = rgb(BLEU_FONCE)
            r.font.name = 'Calibri'
            if sous_titre:
                p2 = tf.add_paragraph()
                r2 = p2.add_run()
                r2.text = sous_titre
                r2.font.size = Pt(14)
                r2.font.italic = True
                r2.font.color.rgb = rgb(GRIS_TEXTE)
            tp = slide.shapes.add_textbox(Inches(12.1), Inches(0.32), Inches(1.0), Inches(0.6))
            ttf = tp.text_frame
            ttf.margin_top = Emu(0)
            pp = ttf.paragraphs[0]
            pp.alignment = PP_ALIGN.RIGHT
            rr = pp.add_run()
            rr.text = numero
            rr.font.size = Pt(14)
            rr.font.bold = True
            rr.font.color.rgb = rgb(ROUGE)
            pic_logo = slide.shapes.add_picture(
                logo_app, Inches(11.6), Inches(0.62), Inches(1.35), Inches(1.35))
            return pic_logo

        def pied(slide, page):
            boite(slide, 0, 7.28, 13.333, 0.22, couleur=BLEU_FONCE)
            texte(slide, 0.6, 7.02, 9.0, 0.3,
                  "© MUKESHABA — " + META['contact'], taille=10, couleur=GRIS)
            texte(slide, 12.3, 7.02, 0.9, 0.3, str(page), taille=10,
                  couleur=GRIS, align=PP_ALIGN.RIGHT)

        def diapo_titre(titre):
            s = prs.slides.add_slide(BLANK)
            fond(s, BLEU_FONCE)
            return s

        # ── Animation : fondu d'entrée successif (After Previous) ──
        def animer(slide, decalage=0):
            from pptx.oxml import parse_xml
            NS = ('xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
                  'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"')
            spids = [sh.shape_id for sh in slide.shapes
                     if sh.shape_type is not None]
            if not spids:
                return
            parts = ['<p:timing %s><p:tnLst><p:par><p:cTn id="1" dur="indefinite" '
                     'restart="never" nodeType="tmRoot"><p:childTnLst>'
                     '<p:seq concurrent="1" nextAc="seek">'
                     '<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
                     '<p:par><p:cTn id="3" fill="hold"><p:stCondLst><p:cond delay="indefinite"/>'
                     '</p:stCondLst><p:childTnLst>' % NS]
            eid = 10
            for i, sid in enumerate(spids):
                ntype = 'clickEffect' if i == 0 else 'afterEffect'
                retard = decalage if i == 0 else (decalage + i * 250)
                eid += 1
                cid, sid_set, sid_eff = eid, eid + 1, eid + 2
                eid += 2
                parts.append(
                    '<p:par><p:cTn id="%d" presetID="10" presetClass="entr" '
                    'presetSubtype="0" fill="hold" grpId="0" nodeType="%s">'
                    '<p:stCondLst><p:cond delay="%d"/></p:stCondLst><p:childTnLst>'
                    '<p:set><p:cBhvr><p:cTn id="%d" dur="1" fill="hold">'
                    '<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
                    '<p:tgtEl><p:spTgt spid="%d"/></p:tgtEl>'
                    '<p:attrNameLst><p:attrName>style.visibility</p:attrName>'
                    '</p:attrNameLst></p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
                    '<p:animEffect transition="in" filter="fade"><p:cBhvr>'
                    '<p:cTn id="%d" dur="600" fill="hold"><p:stCondLst>'
                    '<p:cond delay="0"/></p:stCondLst></p:cTn>'
                    '<p:tgtEl><p:spTgt spid="%d"/></p:tgtEl></p:cBhvr></p:animEffect>'
                    '</p:childTnLst></p:cTn></p:par>'
                    % (cid, ntype, retard, sid_set, sid, sid_eff, sid))
            parts.append('</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:seq>'
                         '</p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>')
            slide._element.append(parse_xml(''.join(parts)))

        # ===========================================================
        # 1. COUVERTURE
        # ===========================================================
        s = diapo_titre("COUVERTURE")
        boite(s, 0, 0, 13.333, 0.16, couleur=ROUGE)
        boite(s, 0, 7.34, 13.333, 0.16, couleur=ROUGE)
        image(s, logo_app, 5.17, 0.85, 3.0, 3.0)
        texte(s, 0.8, 0.6, 2.6, 1.0, META['societe'], taille=20,
              couleur=BLANC, gras=True, align=PP_ALIGN.CENTER)
        texte(s, 0.8, 1.15, 2.6, 1.2, "SOCIÉTÉ D'INGÉNIERIE\n& DE CARTOGRAPHIE",
              taille=11, couleur=(0xBF, 0xD0, 0xE8), align=PP_ALIGN.CENTER)
        texte(s, 9.9, 0.7, 2.9, 1.2, "Version %s" % META['app_version'],
              taille=13, couleur=(0xBF, 0xD0, 0xE8), align=PP_ALIGN.RIGHT)
        texte(s, 9.9, 1.1, 2.9, 1.2, META['date_guide'], taille=13,
              couleur=(0xBF, 0xD0, 0xE8), align=PP_ALIGN.RIGHT)
        texte(s, 1.0, 4.15, 11.33, 1.2, META['app_nom'], taille=88,
              couleur=BLANC, gras=True, align=PP_ALIGN.CENTER)
        texte(s, 1.0, 5.5, 11.33, 0.7, META['slogan'], taille=24,
              couleur=(0xD9, 0xE4, 0xF2), align=PP_ALIGN.CENTER)
        boite(s, 5.42, 6.35, 2.5, 0.05, couleur=ROUGE)
        texte(s, 1.0, 6.6, 11.33, 0.5,
              "Présentation officielle — Guide de l'application", taille=16,
              couleur=(0xBF, 0xD0, 0xE8), align=PP_ALIGN.CENTER)
        image(s, logo_editeur, 6.42, 6.95, 0.5, 0.5)
        animer(s)

        # ===========================================================
        # 2. SOMMAIRE
        # ===========================================================
        s = prs.slides.add_slide(BLANK)
        fond(s, FOND)
        entete(s, "Sommaire", "02")
        sommaire = [
            ("03", "MUKMAP en bref"),
            ("04", "Tableau de bord & indicateurs"),
            ("05", "Carte interactive & fonds de carte"),
            ("06", "Collecte de points sur le terrain"),
            ("07", "Zones de sécurité"),
            ("08", "Itinéraires & suivi terrain"),
            ("09", "Projets & activités"),
            ("10", "Rapports professionnels"),
            ("11", "Adduction d'eau (Water Supply)"),
            ("12", "Agents & contrôle d'accès"),
            ("13", "Mode Avancé (PRO)"),
            ("14", "Guide d'utilisation"),
            ("15", "Accessibilité : mobile, tablette, PC, hors ligne"),
            ("16", "À propos de MUKESHABA"),
        ]
        y = 1.15
        for num, lib in sommaire[:7]:
            boite(s, 1.0, y, 0.55, 0.55, couleur=BLEU_FONCE, arrondi=True)
            texte(s, 1.0, y + 0.1, 0.55, 0.35, num, taille=15,
                  couleur=BLANC, gras=True, align=PP_ALIGN.CENTER)
            texte(s, 1.75, y + 0.09, 4.6, 0.4, lib, taille=15, couleur=GRIS_TEXTE)
            y += 0.75
        y = 1.15
        for num, lib in sommaire[7:]:
            boite(s, 7.0, y, 0.55, 0.55, couleur=ROUGE, arrondi=True)
            texte(s, 7.0, y + 0.1, 0.55, 0.35, num, taille=15,
                  couleur=BLANC, gras=True, align=PP_ALIGN.CENTER)
            texte(s, 7.75, y + 0.09, 4.9, 0.4, lib, taille=15, couleur=GRIS_TEXTE)
            y += 0.75
        pied(s, 2)
        animer(s)

        # ===========================================================
        # 3. MUKMAP EN BREF
        # ===========================================================
        s = prs.slides.add_slide(BLANK)
        fond(s, FOND)
        entete(s, "MUKMAP en bref", "03")
        texte(s, 0.7, 1.3, 11.9, 1.4,
              "MUKMAP est une plateforme SIG professionnelle de collecte, de suivi, de "
              "gestion des risques et de sécurité, éditée par MUKESHABA. Elle organise "
              "toutes les données géographiques par projet puis par activité.",
              taille=16, couleur=GRIS_TEXTE)
        boite(s, 0.7, 2.75, 11.93, 0.02, couleur=BLEU_CLAIR)
        puces(s, 0.7, 3.0, 11.9, 4.0, [
            "Cartographie de points d'intérêt : hôpitaux, écoles, marchés, villages, infrastructures…",
            "Délimitation de zones de sécurité et signalement des zones dangereuses",
            "Levés topographiques : repères géodésiques, altitudes, courbes de niveau, profils",
            "Suivi des réseaux d'adduction d'eau : captages, bornes-fontaines, réservoirs",
            "Mode hors ligne pour travailler sans connexion sur le terrain",
            "Rapports d'activités professionnels en PDF, Word et Excel",
            "Imports/exports SIG : GeoJSON, KML, Shapefile, GPX, DXF, plus 6 langues d'interface",
        ], taille=17)
        pied(s, 3)
        animer(s)

        # ===========================================================
        # 4. TABLEAU DE BORD
        # ===========================================================
        s = prs.slides.add_slide(BLANK)
        fond(s, FOND)
        entete(s, "Tableau de bord & indicateurs", "04")
        puces(s, 0.7, 1.35, 5.4, 5.5, [
            "Vue d'ensemble des points encodés, activités et bénéficiaires",
            "Zones de sécurité : dangereuses, sécurisées, sans information",
            "Évolution des activités, points par province, bénéficiaires par mois",
            "Répartition par projet et par catégorie",
            "Activités récentes et panneau d'administration",
            "Thème clair / sombre et sélecteur de langue",
        ], taille=16)
        image(s, os.path.join(captures, 'dashboard.png'), 6.35, 1.35, 6.6, 4.13, bordure=BLEU_CLAIR)
        pied(s, 4)
        animer(s)

        # ===========================================================
        # 5. CARTE INTERACTIVE
        # ===========================================================
        s = prs.slides.add_slide(BLANK)
        fond(s, FOND)
        entete(s, "Carte interactive & fonds de carte", "05")
        puces(s, 0.7, 1.35, 5.4, 5.5, [
            "Carte multi-fonds : rues, satellite, relief, sombre, topographique",
            "Imagerie aérienne (orthophotos drone) et couches WMS",
            "Recherche de lieu et localisation par coordonnées GPS",
            "Identification d'un élément : table attributaire complète",
            "Outils de dessin et de mesure : point, ligne, polygone, cercle, rectangle",
            "Zoom sur objet, plein écran, navigation fluide",
        ], taille=16)
        image(s, os.path.join(captures, 'carte.png'), 6.35, 1.35, 6.6, 4.13, bordure=BLEU_CLAIR)
        pied(s, 5)
        animer(s)

        # ===========================================================
        # 6. COLLECTE DE POINTS
        # ===========================================================
        s = prs.slides.add_slide(BLANK)
        fond(s, FOND)
        entete(s, "Collecte de points sur le terrain", "06")
        puces(s, 0.7, 1.35, 5.4, 5.5, [
            "Catégories : hôpital, école, église, marché, village, pont, route, incident…",
            "Coordonnées GPS automatiques (cliquer sur la carte ou géolocalisation)",
            "Photos géoréférencées : date, heure et position lues dans l'EXIF",
            "Galerie multimédia : photos, vidéos, audio et PDF",
            "Fiche détaillée, historique d'audit et modifications",
            "Tableau « Tous les points » avec recherche et vues",
        ], taille=16)
        image(s, os.path.join(captures, 'points.png'), 6.35, 1.35, 6.6, 4.13, bordure=BLEU_CLAIR)
        pied(s, 6)
        animer(s)

        # ===========================================================
        # 7. ZONES DE SÉCURITÉ
        # ===========================================================
        s = prs.slides.add_slide(BLANK)
        fond(s, FOND)
        entete(s, "Zones de sécurité", "07")
        puces(s, 0.7, 1.35, 5.4, 5.5, [
            "Déclaration de zones dangereuses, sécurisées ou sans information",
            "Zones ponctuelles (rayon) ou délimitées sur la carte",
            "Motif de déclaration, responsable et date",
            "Alertes visibles sur les itinéraires",
            "Zones dangereuses intégrées dans les rapports",
            "Modification et suppression avec journal d'audit",
        ], taille=16)
        image(s, os.path.join(captures, 'zones.png'), 6.35, 1.35, 6.6, 4.13, bordure=BLEU_CLAIR)
        pied(s, 7)
        animer(s)

        # ===========================================================
        # 8. ITINÉRAIRES
        # ===========================================================
        s = prs.slides.add_slide(BLANK)
        fond(s, FOND)
        entete(s, "Itinéraires & suivi terrain", "08")
        puces(s, 0.7, 1.35, 5.4, 5.5, [
            "Tracer un parcours directement sur la carte",
            "Analyse automatique des zones traversées (dangereuses…)",
            "Liste de mes itinéraires avec détail et date",
            "Alertes de sécurité pendant le trajet",
            "Itinéraires intégrés dans les rapports d'activités",
        ], taille=16)
        image(s, os.path.join(captures, 'itineraire.png'), 6.35, 1.35, 6.6, 4.13, bordure=BLEU_CLAIR)
        pied(s, 8)
        animer(s)

        # ===========================================================
        # 9. PROJETS & ACTIVITÉS
        # ===========================================================
        s = prs.slides.add_slide(BLANK)
        fond(s, FOND)
        entete(s, "Projets & activités", "09")
        puces(s, 0.7, 1.35, 5.4, 5.5, [
            "Organisation par projet puis par activité",
            "Modèles d'activités pour accélérer la saisie terrain",
            "Archivage et réactivation des projets",
            "Suggestions automatiques depuis l'historique",
            "Sélection du projet et de l'activité en cours",
        ], taille=16)
        image(s, os.path.join(captures, 'projets.png'), 6.35, 1.35, 6.6, 4.13, bordure=BLEU_CLAIR)
        pied(s, 9)
        animer(s)

        # ===========================================================
        # 10. RAPPORTS
        # ===========================================================
        s = prs.slides.add_slide(BLANK)
        fond(s, FOND)
        entete(s, "Rapports professionnels", "10")
        puces(s, 0.7, 1.35, 5.4, 5.5, [
            "Assistant de rapport en 6 étapes (période, projet, activités, filtres, sections)",
            "Types : journalier, hebdomadaire, mensuel, annuel, personnalisé",
            "Jusqu'à 15 sections : bénéficiaires, agents, terrain, zones, météo…",
            "Export PDF, Word, Excel et impression",
            "Conditions météorologiques intégrées par activité",
        ], taille=16)
        image(s, os.path.join(captures, 'rapport.png'), 6.35, 1.35, 6.6, 4.13, bordure=BLEU_CLAIR)
        pied(s, 10)
        animer(s)

        # ===========================================================
        # 11. ADDUCTION D'EAU
        # ===========================================================
        s = prs.slides.add_slide(BLANK)
        fond(s, FOND)
        entete(s, "Adduction d'eau (Water Supply)", "11")
        puces(s, 0.7, 1.35, 5.4, 5.5, [
            "Module complet de collecte : sources, bornes-fontaines, villages, réservoirs",
            "Données techniques : débit, profondeur, qualité de l'eau (pH, turbidité…)",
            "Dessin des conduites et tracé du réseau",
            "Tableau de bord d'indicateurs et suivi du projet",
            "Export PDF / Excel des ouvrages et profils",
        ], taille=16)
        image(s, os.path.join(captures, 'adduction.png'), 6.35, 1.35, 6.6, 4.13, bordure=BLEU_CLAIR)
        pied(s, 11)
        animer(s)

        # ===========================================================
        # 12. AGENTS
        # ===========================================================
        s = prs.slides.add_slide(BLANK)
        fond(s, FOND)
        entete(s, "Agents & contrôle d'accès", "12")
        puces(s, 0.7, 1.35, 5.4, 5.5, [
            "Gestion des comptes agents de terrain",
            "Blocage et déblocage en un clic",
            "Profil avec fonction, mission et photo d'identité",
            "Journal d'audit complet (action, agent, IP, date)",
            "Mot de passe oublié géré par l'administrateur",
        ], taille=16)
        image(s, os.path.join(captures, 'agents.png'), 6.35, 1.35, 6.6, 4.13, bordure=BLEU_CLAIR)
        pied(s, 12)
        animer(s)

        # ===========================================================
        # 13. MODE AVANCÉ
        # ===========================================================
        s = prs.slides.add_slide(BLANK)
        fond(s, FOND)
        entete(s, "Mode Avancé (PRO)", "13")
        puces(s, 0.7, 1.35, 5.4, 5.5, [
            "Code d'accès temporaire ou permanent, réservé à l'administrateur principal",
            "Fonds de carte personnalisés et imagerie aérienne",
            "Couches WMS superposables",
            "Outils de terrain avancés : topographie, infrastructures, réseau eau",
            "Le mode avancé est détaillé dans le guide d'utilisation",
        ], taille=16)
        image(s, os.path.join(captures, 'mode_avance.png'), 6.35, 1.35, 6.6, 4.13, bordure=BLEU_CLAIR)
        pied(s, 13)
        animer(s)

        # ===========================================================
        # 14. GUIDE D'UTILISATION
        # ===========================================================
        s = prs.slides.add_slide(BLANK)
        fond(s, FOND)
        entete(s, "Guide d'utilisation", "14")
        puces(s, 0.7, 1.35, 5.4, 5.5, [
            "Guide complet officiel en PDF et Word (20 chapitres + annexes)",
            "Accessible dans l'espace Super Admin",
            "Glossaire, outils/icônes, raccourcis, FAQ, procédures types",
            "Version affichée : MUKMAP 1.0",
            "Édité et publié par MUKESHABA",
        ], taille=16)
        image(s, os.path.join(captures, 'guide.png'), 6.35, 1.35, 6.6, 4.13, bordure=BLEU_CLAIR)
        pied(s, 14)
        animer(s)

        # ===========================================================
        # 15. ACCESSIBILITÉ / HORS LIGNE
        # ===========================================================
        s = prs.slides.add_slide(BLANK)
        fond(s, FOND)
        entete(s, "Accessibilité : mobile, tablette, PC, hors ligne", "15")
        puces(s, 0.7, 1.35, 5.4, 5.5, [
            "Application installable (PWA) sur téléphone, tablette et ordinateur",
            "Interface adaptative : tiroir plein écran, panneaux flottants, FAB",
            "Mode hors ligne : téléchargement des zones, collecte puis synchronisation",
            "Météo automatique avec cache local et badge temps réel",
            "6 langues : français, anglais, swahili, lingala, portugais, chinois",
        ], taille=16)
        image(s, os.path.join(captures, 'connexion.png'), 6.35, 1.35, 6.6, 4.13, bordure=BLEU_CLAIR)
        pied(s, 15)
        animer(s)

        # ===========================================================
        # 16. À PROPOS / CONTACT
        # ===========================================================
        s = diapo_titre("CONTACT")
        boite(s, 0, 0, 13.333, 0.16, couleur=ROUGE)
        boite(s, 0, 7.34, 13.333, 0.16, couleur=ROUGE)
        image(s, logo_app, 5.42, 0.9, 2.5, 2.5)
        texte(s, 1.0, 3.6, 11.33, 0.8, "À propos de MUKESHABA", taille=36,
              couleur=BLANC, gras=True, align=PP_ALIGN.CENTER)
        texte(s, 1.6, 4.5, 10.13, 1.1,
              "MUKESHABA est une entreprise de conception et de développement de logiciels. "
              "MUKMAP est développée par Ir. Yagirwa Gedeon — Consultant SIG & Développeur Web.",
              taille=16, couleur=(0xD9, 0xE4, 0xF2), align=PP_ALIGN.CENTER)
        boite(s, 4.42, 5.85, 4.5, 0.9, couleur=BLEU_CLAIR, arrondi=True)
        texte(s, 4.42, 6.0, 4.5, 0.6, META['contact'], taille=14,
              couleur=BLANC, gras=True, align=PP_ALIGN.CENTER)
        image(s, logo_editeur, 6.42, 6.9, 0.5, 0.5)
        animer(s)

        prs.save(chemin)
        return chemin